import pymupdf
import pymupdf4llm
from pptx import Presentation
from docx import Document
import io


def extract_pdf(contents: bytes) -> str:
    doc = pymupdf.open(stream=contents, filetype="pdf")
    return pymupdf4llm.to_text(doc)


def extract_pptx(contents: bytes) -> str:
    prs = Presentation(io.BytesIO(contents))
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text_runs.append(para.text)
    return "\n".join(text_runs)


def extract_docx(contents: bytes) -> str:
    doc = Document(io.BytesIO(contents))
    return "\n".join(p.text for p in doc.paragraphs)


def extract_txt(contents: bytes) -> str:
    return contents.decode("utf-8", errors="replace")


EXTRACTORS = {
    "pdf": extract_pdf,
    "pptx": extract_pptx,
    "docx": extract_docx,
    "txt": extract_txt,
}

def extract_text(filename: str, contents: bytes) -> str:
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext == "pdf":
        return extract_pdf(contents)
    elif ext == "pptx":
        return extract_pptx(contents)
    elif ext == "docx":
        return extract_docx(contents)
    elif ext == "txt":
        return extract_txt(contents)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")